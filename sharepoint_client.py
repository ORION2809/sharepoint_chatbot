"""SharePoint client using Microsoft Graph API (delegated auth with token cache)."""

from __future__ import annotations

import io
import json
import logging
from pathlib import Path
from typing import Any

import httpx
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

import config

logger = logging.getLogger(__name__)

GRAPH_BASE = "https://graph.microsoft.com/v1.0"
GRAPH_SCOPES = ["Sites.Read.All", "Files.Read.All"]
TOKEN_CACHE_FILE = Path(__file__).parent / ".token_cache.json"


class SharePointClient:
    def __init__(self) -> None:
        self._token: str | None = None
        self._refresh_token: str | None = None
        self._site_id: str | None = None
        # Reusable HTTP clients with connection pooling
        self._http: httpx.Client = httpx.Client(timeout=30)
        self._http_dl: httpx.Client = httpx.Client(timeout=60, follow_redirects=True)
        self._load_cache()

    def _load_cache(self) -> None:
        logger.info("Token cache path: %s (exists=%s)", TOKEN_CACHE_FILE, TOKEN_CACHE_FILE.exists())
        if not TOKEN_CACHE_FILE.exists():
            return
        try:
            data = json.loads(TOKEN_CACHE_FILE.read_text())
            self._token = data.get("access_token") or None
            self._refresh_token = data.get("refresh_token") or None
            logger.info("Loaded token from cache (token=%s, refresh=%s)",
                        bool(self._token), bool(self._refresh_token))
        except Exception as exc:
            logger.warning("Failed to load token cache: %s", exc)

    def _save_cache(self) -> None:
        TOKEN_CACHE_FILE.write_text(json.dumps({
            "access_token": self._token or "",
            "refresh_token": self._refresh_token or "",
        }))

    # ── Auth ────────────────────────────────────────────────────────────
    def _get_token(self) -> str:
        if self._token:
            return self._token

        if self._refresh_token:
            self._do_refresh()
            if self._token:
                return self._token

        raise RuntimeError(
            "Not authenticated! Run 'python login.py' first to sign in."
        )

    def _do_refresh(self) -> None:
        """Refresh the access token using the stored refresh token."""
        token_url = f"https://login.microsoftonline.com/{config.AZURE_TENANT_ID}/oauth2/v2.0/token"
        scope_str = " ".join(GRAPH_SCOPES) + " offline_access"
        data = {
            "client_id": config.AZURE_CLIENT_ID,
            "client_secret": config.AZURE_CLIENT_SECRET,
            "refresh_token": self._refresh_token,
            "grant_type": "refresh_token",
            "scope": scope_str,
        }
        try:
            resp = self._http.post(token_url, data=data)
            result = resp.json()
            if "access_token" in result:
                self._token = result["access_token"]
                self._refresh_token = result.get("refresh_token", self._refresh_token)
                self._save_cache()
                logger.info("Token refreshed successfully.")
            else:
                logger.warning("Token refresh failed: %s", result.get("error_description", "unknown"))
                self._token = None
        except Exception as exc:
            logger.warning("Token refresh error: %s", exc)

    def _headers(self) -> dict[str, str]:
        token = self._token or self._get_token()
        return {"Authorization": f"Bearer {token}", "Accept": "application/json"}

    def _get(self, url: str) -> dict[str, Any]:
        resp = self._http.get(url, headers=self._headers())
        if resp.status_code == 401:
            # Token expired — refresh and retry
            self._token = None
            self._do_refresh()
            resp = self._http.get(url, headers=self._headers())
        resp.raise_for_status()
        return resp.json()

    def _download(self, url: str) -> bytes:
        resp = self._http_dl.get(url, headers=self._headers())
        if resp.status_code == 401:
            self._token = None
            self._do_refresh()
            resp = self._http_dl.get(url, headers=self._headers())
        resp.raise_for_status()
        return resp.content

    def _get_paged(self, url: str, max_pages: int = 50) -> list[dict[str, Any]]:
        """GET with automatic @odata.nextLink pagination."""
        all_items: list[dict[str, Any]] = []
        page = 0
        while url and page < max_pages:
            data = self._get(url)
            all_items.extend(data.get("value", []))
            url = data.get("@odata.nextLink")  # type: ignore[assignment]
            page += 1
        return all_items

    # ── Site discovery ──────────────────────────────────────────────────
    def get_site_id(self) -> str:
        if self._site_id:
            return self._site_id

        hostname = config.SHAREPOINT_HOSTNAME
        site_path = config.SHAREPOINT_SITE_PATH

        if site_path:
            url = f"{GRAPH_BASE}/sites/{hostname}:/{site_path}"
        else:
            url = f"{GRAPH_BASE}/sites/{hostname}:/"

        data = self._get(url)
        self._site_id = data["id"]
        logger.info("Resolved site id: %s", self._site_id)
        return self._site_id

    # ── List drives / files ─────────────────────────────────────────────
    def list_drives(self) -> list[dict[str, Any]]:
        site_id = self.get_site_id()
        data = self._get(f"{GRAPH_BASE}/sites/{site_id}/drives")
        return data.get("value", [])

    def list_files(self, drive_id: str, folder_path: str = "") -> list[dict[str, Any]]:
        if folder_path:
            url = f"{GRAPH_BASE}/drives/{drive_id}/root:/{folder_path}:/children?$top=200"
        else:
            url = f"{GRAPH_BASE}/drives/{drive_id}/root/children?$top=200"
        return self._get_paged(url)

    def list_all_files_recursive(
        self, drive_id: str, folder_path: str = "", depth: int = 3
    ) -> list[dict[str, Any]]:
        if depth <= 0:
            return []
        items = self.list_files(drive_id, folder_path)
        result: list[dict[str, Any]] = []
        for item in items:
            if "file" in item:
                result.append(item)
            elif "folder" in item and depth > 1:
                child_path = f"{folder_path}/{item['name']}" if folder_path else item["name"]
                result.extend(self.list_all_files_recursive(drive_id, child_path, depth - 1))
        return result

    # ── All-sites scanning (Priority 4) ──────────────────────────────
    def list_all_sites(self) -> list[dict[str, Any]]:
        """Return every SharePoint site the signed-in user can access."""
        url = f"{GRAPH_BASE}/sites?search=*&$top=100"
        return self._get_paged(url, max_pages=20)

    def list_site_drives(self, site_id: str) -> list[dict[str, Any]]:
        """Return drives (document libraries) for a given site."""
        data = self._get(f"{GRAPH_BASE}/sites/{site_id}/drives")
        return data.get("value", [])

    def list_all_accessible_files(
        self, depth: int = 3, on_progress: Any = None,
    ) -> list[dict[str, Any]]:
        """Scan ALL SharePoint sites + drives the user can access.

        *on_progress* is an optional callback(msg: str) for status updates.
        """
        sites = self.list_all_sites()
        # Always include the configured site even if the search misses it
        try:
            cfg_site_id = self.get_site_id()
            if not any(s["id"] == cfg_site_id for s in sites):
                cfg_data = self._get(f"{GRAPH_BASE}/sites/{cfg_site_id}")
                sites.insert(0, cfg_data)
        except Exception:
            pass

        seen_ids: set[str] = set()
        all_files: list[dict[str, Any]] = []

        for idx, site in enumerate(sites):
            sid = site["id"]
            if sid in seen_ids:
                continue
            seen_ids.add(sid)
            site_name = site.get("displayName", sid)

            if on_progress:
                on_progress(f"Scanning site {idx + 1}/{len(sites)}: {site_name}")

            try:
                drives = self.list_site_drives(sid)
                for drive in drives:
                    files = self.list_all_files_recursive(
                        drive["id"], "", depth=depth,
                    )
                    for f in files:
                        f["_site_name"] = site_name
                        f["_site_id"] = sid
                    all_files.extend(files)
            except Exception as exc:
                logger.warning("Failed to scan site '%s': %s", site_name, exc)

        logger.info("All-sites scan complete: %d files across %d sites",
                    len(all_files), len(seen_ids))
        return all_files

    # ── Search ──────────────────────────────────────────────────────────
    def search_files(self, query: str) -> list[dict[str, Any]]:
        site_id = self.get_site_id()
        # Sanitize query: remove quotes and special chars that break OData
        safe_query = query.replace("'", "").replace('"', "").replace("?", "").strip()
        # Use first few keywords to keep the search broad
        keywords = safe_query.split()[:5]
        search_term = " ".join(keywords)
        url = f"{GRAPH_BASE}/sites/{site_id}/drive/root/search(q='{search_term}')"
        try:
            data = self._get(url)
            return data.get("value", [])
        except Exception as exc:
            logger.warning("Search failed, falling back to drive listing: %s", exc)
            return []

    # ── Download & extract text ─────────────────────────────────────────
    def download_file_content(self, download_url: str) -> bytes:
        return self._download(download_url)

    def download_raw(self, file_item: dict[str, Any]) -> bytes | None:
        """Download raw bytes for a file item. Returns None on failure."""
        name = file_item.get("name", "")
        download_url = file_item.get("@microsoft.graph.downloadUrl", "")

        if not download_url:
            item_id = file_item.get("id", "")
            parent_ref = file_item.get("parentReference", {})
            drive_id = parent_ref.get("driveId", "")
            if drive_id and item_id:
                url = f"{GRAPH_BASE}/drives/{drive_id}/items/{item_id}/content"
                try:
                    return self._download(url)
                except Exception as exc:
                    logger.warning("Failed to download %s: %s", name, exc)
                    return None
            return None

        try:
            return self._download(download_url)
        except Exception as exc:
            logger.warning("Failed to download %s: %s", name, exc)
            return None

    def extract_text(self, file_item: dict[str, Any], use_vision: bool = False) -> str:
        """Extract text from a file. If use_vision=False, images/videos get metadata only."""
        name = file_item.get("name", "")
        raw = self.download_raw(file_item)
        if raw is None:
            return ""
        return self._parse_bytes(name, raw, use_vision=use_vision)

    @staticmethod
    def _parse_bytes(name: str, raw: bytes, *, use_vision: bool = False) -> str:
        lower = name.lower()
        if lower.endswith(".txt") or lower.endswith(".csv") or lower.endswith(".md"):
            return raw.decode("utf-8", errors="replace")
        if lower.endswith(".pdf"):
            reader = PdfReader(io.BytesIO(raw))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        if lower.endswith(".docx"):
            doc = DocxDocument(io.BytesIO(raw))
            return "\n".join(p.text for p in doc.paragraphs)
        if lower.endswith((".html", ".htm")):
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(raw, "html.parser")
            return soup.get_text(separator="\n", strip=True)
        if lower.endswith(".xlsx"):
            wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            parts: list[str] = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                parts.append(f"--- Sheet: {sheet} ---")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        parts.append("\t".join(cells))
            wb.close()
            return "\n".join(parts)
        if lower.endswith(".pptx"):
            prs = Presentation(io.BytesIO(raw))
            parts = []
            for idx, slide in enumerate(prs.slides, 1):
                parts.append(f"--- Slide {idx} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text for cell in row.cells]
                            parts.append("\t".join(cells))
            return "\n".join(parts)
        if lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff")):
            if use_vision:
                from vision_client import describe_image
                return describe_image(name, raw)
            return f"[Image file: {name} | size: {len(raw)} bytes]"
        if lower.endswith((".mp4", ".avi", ".mov", ".wmv", ".mkv", ".webm")):
            if use_vision:
                from vision_client import describe_video
                return describe_video(name, raw)
            return f"[Video file: {name} | size: {len(raw)} bytes]"
        # Fallback: try as plain text
        try:
            return raw.decode("utf-8", errors="replace")
        except Exception:
            return f"[Could not extract text from {name}]"
